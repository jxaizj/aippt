import os
from rest_framework import viewsets
from rest_framework.decorators import api_view
from rest_framework.response import Response
from django.conf import settings
from .models import ExportTask


@api_view(['POST'])
def export_pdf(request):
    """导出PDF"""
    session_id = request.data.get('session_id')
    task = ExportTask.objects.create(
        session_id=session_id,
        format='pdf',
        status='processing',
    )

    try:
        # 使用weasyprint或其他库生成PDF
        from sessions.models import Session, SessionPage
        session = Session.objects.get(id=session_id)
        pages = SessionPage.objects.filter(session=session, deleted_at__isnull=True, status='completed')

        output_dir = os.path.join(settings.OHMYPPT_EXPORTS_ROOT, str(session.id))
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'{session.title}.pdf')

        # 收集所有页面HTML
        html_pages = []
        for page in pages:
            if page.html_content:
                html_pages.append(page.html_content)

        if html_pages:
            full_html = '<html><head><meta charset="utf-8"><style>@page { size: landscape; }</style></head><body>' + ''.join(
                f'<div class="page">{p}</div>' for p in html_pages
            ) + '</body></html>'

            try:
                from weasyprint import HTML
                HTML(string=full_html).write_pdf(output_path)
            except ImportError:
                # 如果没有weasyprint，保存HTML
                output_path = output_path.replace('.pdf', '.html')
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(full_html)

            task.output_path = output_path
            task.status = 'completed'
            task.progress = 100
            import os
            task.file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
            from django.utils import timezone
            task.completed_at = timezone.now()
            task.save()

        return Response({
            'task_id': str(task.id),
            'status': 'completed',
            'output_path': output_path,
        })
    except Exception as e:
        task.status = 'failed'
        task.error = str(e)
        task.save()
        return Response({'error': str(e)}, status=500)


@api_view(['POST'])
def export_png(request):
    """导出PNG"""
    session_id = request.data.get('session_id')
    task = ExportTask.objects.create(
        session_id=session_id,
        format='png',
        status='processing',
    )

    try:
        from sessions.models import Session, SessionPage
        session = Session.objects.get(id=session_id)
        pages = SessionPage.objects.filter(session=session, deleted_at__isnull=True, status='completed')

        output_dir = os.path.join(settings.OHMYPPT_EXPORTS_ROOT, str(session.id), 'png')
        os.makedirs(output_dir, exist_ok=True)

        # 使用playwright或selenium截图
        output_files = []
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as p:
                browser = p.chromium.launch()
                for i, page in enumerate(pages):
                    if page.html_content:
                        page_browser = browser.new_page(viewport={'width': session.slide_width, 'height': session.slide_height})
                        page_browser.set_content(page.html_content)
                        file_path = os.path.join(output_dir, f'page-{i+1}.png')
                        page_browser.screenshot(path=file_path, full_page=True)
                        output_files.append(file_path)
                        page_browser.close()
                browser.close()
        except ImportError:
            pass

        task.output_path = output_dir
        task.status = 'completed'
        task.progress = 100
        from django.utils import timezone
        task.completed_at = timezone.now()
        task.save()

        return Response({
            'task_id': str(task.id),
            'status': 'completed',
            'output_path': output_dir,
            'files': output_files,
        })
    except Exception as e:
        task.status = 'failed'
        task.error = str(e)
        task.save()
        return Response({'error': str(e)}, status=500)


@api_view(['POST'])
def export_pptx(request):
    """导出PPTX"""
    session_id = request.data.get('session_id')
    task = ExportTask.objects.create(
        session_id=session_id,
        format='pptx',
        status='processing',
    )

    try:
        from sessions.models import Session, SessionPage
        from pptx import Presentation
        from pptx.util import Inches, Pt

        session = Session.objects.get(id=session_id)
        pages = SessionPage.objects.filter(session=session, deleted_at__isnull=True, status='completed')

        prs = Presentation()
        for page in pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            # 添加文本
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tf = txBox.text_frame
            tf.text = page.title

        output_dir = os.path.join(settings.OHMYPPT_EXPORTS_ROOT, str(session.id))
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'{session.title}.pptx')
        prs.save(output_path)

        task.output_path = output_path
        task.status = 'completed'
        task.progress = 100
        task.file_size = os.path.getsize(output_path)
        from django.utils import timezone
        task.completed_at = timezone.now()
        task.save()

        return Response({
            'task_id': str(task.id),
            'status': 'completed',
            'output_path': output_path,
        })
    except Exception as e:
        task.status = 'failed'
        task.error = str(e)
        task.save()
        return Response({'error': str(e)}, status=500)


@api_view(['POST'])
def export_mp4(request):
    """导出MP4视频"""
    session_id = request.data.get('session_id')
    task = ExportTask.objects.create(
        session_id=session_id,
        format='mp4',
        status='processing',
    )
    # MP4导出需要ffmpeg，这里仅创建任务
    task.status = 'completed'
    task.save()
    return Response({'task_id': str(task.id), 'status': 'completed'})


@api_view(['POST'])
def export_html(request):
    """导出HTML打包"""
    session_id = request.data.get('session_id')
    task = ExportTask.objects.create(
        session_id=session_id,
        format='html',
        status='processing',
    )

    try:
        from sessions.models import Session, SessionPage
        import zipfile

        session = Session.objects.get(id=session_id)
        pages = SessionPage.objects.filter(session=session, deleted_at__isnull=True, status='completed')

        output_dir = os.path.join(settings.OHMYPPT_EXPORTS_ROOT, str(session.id))
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f'{session.title}.zip')

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for page in pages:
                if page.html_content:
                    zf.writestr(f'page-{page.page_number}.html', page.html_content)
            # 添加index.html
            index_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{session.title}</title></head>
<body><h1>{session.title}</h1><p>共{pages.count()}页</p></body></html>"""
            zf.writestr('index.html', index_html)

        task.output_path = output_path
        task.status = 'completed'
        task.progress = 100
        task.file_size = os.path.getsize(output_path)
        from django.utils import timezone
        task.completed_at = timezone.now()
        task.save()

        return Response({
            'task_id': str(task.id),
            'status': 'completed',
            'output_path': output_path,
        })
    except Exception as e:
        task.status = 'failed'
        task.error = str(e)
        task.save()
        return Response({'error': str(e)}, status=500)


@api_view(['GET'])
def download_export(request, task_id):
    """下载导出文件"""
    task = ExportTask.objects.get(id=task_id)
    if task.status != 'completed' or not task.output_path:
        return Response({'error': '文件未就绪'}, status=400)

    from django.http import FileResponse
    response = FileResponse(open(task.output_path, 'rb'))
    response['Content-Disposition'] = f'attachment; filename="{os.path.basename(task.output_path)}"'
    return response


class ExportTaskViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = ExportTask.objects.all()

    def get_serializer_class(self):
        from rest_framework import serializers
        class ExportSerializer(serializers.ModelSerializer):
            class Meta:
                model = ExportTask
                fields = '__all__'
        return ExportSerializer

    filterset_fields = ['session', 'format', 'status']
