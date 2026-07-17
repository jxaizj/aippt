import threading
from rest_framework import viewsets, status, serializers
from rest_framework.decorators import api_view, action
from rest_framework.response import Response
from sessions.models import Session, SessionPage, Message, GenerationRun, GenerationPage, SessionOperation, SpeechScript, ModelConfig
from .models import AIConversation
from .utils import get_active_provider, get_openai_client, save_uploaded_file, safe_json_parse


def _make_progress_sender(session_id):
    """创建进度推送函数，广播到 WebSocket 群组。"""
    from channels.layers import get_channel_layer
    from asgiref.sync import async_to_sync

    channel_layer = get_channel_layer()
    group = f'generation_{session_id}'

    def send(event):
        try:
            async_to_sync(channel_layer.group_send)(group, event)
        except Exception:
            pass
    return send


def _run_in_thread(target, *args, **kwargs):
    """在后台守护线程中运行目标函数，并正确管理 Django 数据库连接。"""
    from django.db import close_old_connections

    def _worker():
        close_old_connections()
        try:
            target(*args, **kwargs)
        finally:
            close_old_connections()

    t = threading.Thread(target=_worker, daemon=True)
    t.start()


class GeneratePptSerializer(serializers.Serializer):
    session_id = serializers.UUIDField()
    topic = serializers.CharField(allow_blank=True, default='')
    page_count = serializers.IntegerField(default=10)
    style_id = serializers.UUIDField(required=False, allow_null=True)
    detail = serializers.CharField(allow_blank=True, default='')
    model_config_id = serializers.UUIDField(required=False, allow_null=True)


class ChatEditSerializer(serializers.Serializer):
    session_id = serializers.UUIDField()
    message = serializers.CharField(allow_blank=True, default='')
    page_id = serializers.CharField(required=False, allow_null=True)
    selector = serializers.CharField(required=False, allow_null=True)


@api_view(['POST'])
def generate_ppt(request):
    """AI生成PPT"""
    serializer = GeneratePptSerializer(data=request.data)
    serializer.is_valid(raise_exception=True)
    data = serializer.validated_data

    session_id = data['session_id']
    topic = data['topic']
    page_count = data['page_count']
    style_id = data.get('style_id')
    detail = data.get('detail', '')
    model_config_id = data.get('model_config_id')

    try:
        session = Session.objects.get(id=session_id)
    except Session.DoesNotExist:
        return Response({'error': '会话不存在'}, status=404)

    # 解析指定的模型配置，回退到激活配置
    model_config = None
    if model_config_id:
        model_config = ModelConfig.objects.filter(id=model_config_id, active=True).first()
    if not model_config:
        model_config = ModelConfig.objects.filter(active=True).first()

    run = GenerationRun.objects.create(
        session=session,
        mode='generate',
        total_pages=page_count,
        status='running',
        model_config=model_config,
    )

    for i in range(1, page_count + 1):
        GenerationPage.objects.create(
            run=run,
            session=session,
            page_id=f'page-{i}',
            page_number=i,
            title=f'第{i}页',
            status='pending'
        )

    conversation = AIConversation.objects.create(
        session=session,
        task_type='generate',
        messages=[{'role': 'user', 'content': f'生成关于"{topic}"的PPT，共{page_count}页。{detail}'}]
    )

    operation = SessionOperation.objects.create(
        session=session,
        type='generate',
        scope='session',
        prompt=topic,
        status='committing',
    )

    # 在后台线程执行生成，避免阻塞 HTTP 请求（生成需多次调用 LLM，耗时较长）
    progress = _make_progress_sender(session_id)
    from .tasks import generate_ppt_task
    _run_in_thread(
        generate_ppt_task,
        str(session_id), str(run.id), str(conversation.id),
        topic, page_count, style_id, detail, model_config,
        progress_callback=progress,
    )

    return Response({
        'session_id': str(session_id),
        'run_id': str(run.id),
        'conversation_id': str(conversation.id),
        'operation_id': str(operation.id),
        'status': 'running',
    })


@api_view(['POST'])
def chat_edit(request):
    """对话式修改PPT"""
    serializer = ChatEditSerializer(data=request.data)
    serializer.is_valid(raise_exception=True)
    data = serializer.validated_data

    session_id = data['session_id']
    message = data['message']
    page_id = data.get('page_id')
    selector = data.get('selector')

    try:
        session = Session.objects.get(id=session_id)
    except Session.DoesNotExist:
        return Response({'error': '会话不存在'}, status=404)

    Message.objects.create(
        session=session,
        role='user',
        content=message,
        type='text',
        chat_scope='page' if page_id else 'main',
        page_id=page_id,
        selector=selector,
    )

    conversation = AIConversation.objects.create(
        session=session,
        task_type='edit',
        messages=[{'role': 'user', 'content': message}]
    )

    operation = SessionOperation.objects.create(
        session=session,
        type='edit',
        scope='page' if page_id else 'session',
        prompt=message,
        status='committing',
    )

    from .tasks import chat_edit_task
    chat_edit_task(str(session_id), str(conversation.id), str(operation.id),
                   message, page_id, selector)

    return Response({
        'session_id': str(session_id),
        'conversation_id': str(conversation.id),
        'operation_id': str(operation.id),
        'status': 'running'
    })


@api_view(['POST'])
def generate_speech(request):
    """生成演讲稿"""
    session_id = request.data.get('session_id')
    page_id = request.data.get('page_id')
    style = request.data.get('style', 'formal')

    from .tasks import generate_speech_task
    result = generate_speech_task(str(session_id), page_id, style)

    return Response({
        'result': result,
        'status': 'completed'
    })


@api_view(['POST'])
def generate_outline(request):
    """生成大纲"""
    topic = request.data.get('topic', '')
    page_count = request.data.get('page_count', 10)
    detail = request.data.get('detail', '')

    provider = get_active_provider()
    if not provider:
        return Response({'error': '未配置AI模型'}, status=400)

    try:
        client = get_openai_client(provider)

        prompt = f"""请为以下主题生成一份详细的PPT大纲，共{page_count}页。
主题：{topic}
补充说明：{detail}

请以JSON数组格式返回，每页包含title(标题)、content(内容描述)、layout(布局建议)。"""

        response = client.chat.completions.create(
            model='gpt-4',
            messages=[
                {'role': 'system', 'content': '你是一个专业的PPT大纲设计师，擅长将主题分解为结构清晰的演示大纲。只返回JSON格式。'},
                {'role': 'user', 'content': prompt}
            ],
            temperature=0.7,
        )

        content = response.choices[0].message.content
        outline = safe_json_parse(content, r'\[.*\]', [])

        return Response({'outline': outline})
    except Exception as e:
        return Response({'error': str(e)}, status=500)


@api_view(['POST'])
def upload_document(request):
    """上传文档解析"""
    file = request.FILES.get('file')
    if not file:
        return Response({'error': '未上传文件'}, status=400)

    file_path = save_uploaded_file(file, 'documents')

    content = ''
    file_ext = os.path.splitext(file.name)[1].lower()
    try:
        if file_ext in ('.txt', '.md'):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_ext == '.docx':
            import docx
            doc = docx.Document(file_path)
            content = '\n'.join([p.text for p in doc.paragraphs])
        elif file_ext == '.csv':
            import csv
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                content = '\n'.join([','.join(row) for row in reader])
        elif file_ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c else '' for c in row])
            content = '\n'.join([' | '.join(row) for row in rows])
    except Exception as e:
        return Response({'error': f'文档解析失败: {str(e)}'}, status=500)

    return Response({'content': content, 'file_path': file_path, 'file_name': file.name})


@api_view(['POST'])
def import_pptx(request):
    """导入PPTX文件"""
    file = request.FILES.get('file')
    if not file:
        return Response({'error': '未上传文件'}, status=400)

    file_path = save_uploaded_file(file, 'pptx_imports')

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            slide_data = {
                'page_number': i + 1,
                'title': f'第{i+1}页',
                'texts': [],
                'layout': slide.slide_layout.name if slide.slide_layout else '',
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_data['texts'].append(shape.text_frame.text)
            pages.append(slide_data)

        session = Session.objects.create(
            title=file.name.replace('.pptx', ''),
            topic=f'导入: {file.name}',
            page_count=len(pages),
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            status='active',
        )

        for page_data in pages:
            SessionPage.objects.create(
                session=session,
                file_slug=f'page-{page_data["page_number"]}',
                page_number=page_data['page_number'],
                title=page_data['title'],
                html_path=f'/sessions/{session.id}/pages/page-{page_data["page_number"]}.html',
                status='completed',
                html_content=f'<div class="slide"><h1>{page_data["title"]}</h1>{"".join(f"<p>{t}</p>" for t in page_data["texts"])}</div>'
            )

        return Response({
            'session_id': str(session.id),
            'pages': pages,
            'page_count': len(pages),
        })
    except Exception as e:
        return Response({'error': f'PPTX解析失败: {str(e)}'}, status=500)


@api_view(['POST'])
def analyze_image(request):
    """图片识别生成风格与大纲"""
    file = request.FILES.get('file')
    if not file:
        return Response({'error': '未上传图片'}, status=400)

    file_path = save_uploaded_file(file, 'temp_images')

    import base64
    with open(file_path, 'rb') as f:
        img_base64 = base64.b64encode(f.read()).decode()

    provider = get_active_provider()
    if not provider:
        return Response({'error': '未配置AI模型'}, status=400)

    try:
        client = get_openai_client(provider)

        response = client.chat.completions.create(
            model='gpt-4o',
            messages=[
                {'role': 'system', 'content': '你是一个视觉设计专家，擅长分析设计稿/截图的视觉特征并提取设计风格。'},
                {'role': 'user', 'content': [
                    {'type': 'text', 'text': '请分析这张图片的视觉风格特征（配色、排版、字体风格等），并基于此风格生成一份PPT大纲。以JSON格式返回：{"style": {"colors": [], "fonts": [], "description": ""}, "outline": [{"title": "", "content": ""}]}'},
                    {'type': 'image_url', 'image_url': {'url': f'data:image/png;base64,{img_base64}'}}
                ]},
            ],
            temperature=0.7,
        )

        content = response.choices[0].message.content
        result = safe_json_parse(content, r'\{.*\}', {'style': {}, 'outline': []})

        return Response(result)
    except Exception as e:
        return Response({'error': str(e)}, status=500)


class AIConversationViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = AIConversation.objects.all()

    def get_serializer_class(self):
        from rest_framework import serializers
        class ConvSerializer(serializers.ModelSerializer):
            class Meta:
                model = AIConversation
                fields = '__all__'
        return ConvSerializer
